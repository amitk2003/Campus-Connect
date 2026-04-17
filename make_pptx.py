"""
CampusConnect Presentation → PowerPoint Generator
Generates presentation.pptx matching the HTML slide content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree

# ─── Colours ────────────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0F, 0x0F, 0x1A)   # main dark bg
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
C_INDIGO    = RGBColor(0x4F, 0x46, 0xE5)
C_PINK      = RGBColor(0xEC, 0x48, 0x99)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x94, 0xA3, 0xB8)
C_GREEN     = RGBColor(0x4A, 0xDE, 0x80)
C_RED       = RGBColor(0xF8, 0x71, 0x71)
C_BLUE      = RGBColor(0x60, 0xA5, 0xFA)
C_ORANGE    = RGBColor(0xFB, 0x92, 0x3C)
C_VIOLET    = RGBColor(0xC4, 0xB5, 0xFD)
C_YELLOW    = RGBColor(0xFA, 0xCC, 0x15)
C_CARD      = RGBColor(0x1E, 0x1E, 0x35)

# ─── Slide dimensions (16:9 widescreen) ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, word_wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_chip(slide, text, left, top):
    """Small pill-label."""
    w, h = Inches(3.5), Inches(0.35)
    r = add_rect(slide, left, top, w, h, fill_color=RGBColor(0x20,0x20,0x38),
                 line_color=C_PURPLE, line_width=Pt(1))
    r.line.width = Pt(1)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = C_PURPLE


def h_center(width):
    """Return left offset to horizontally centre an element of `width`."""
    return (W - width) / 2


def card(slide, left, top, width, height, title, body, icon="",
         title_color=C_WHITE, body_color=C_MUTED, border=C_PURPLE):
    bg = add_rect(slide, left, top, width, height,
                  fill_color=RGBColor(0x1A,0x1A,0x30),
                  line_color=border, line_width=Pt(1))
    pad = Inches(0.2)
    y = top + pad
    if icon:
        add_text(slide, icon, left+pad, y, width-2*pad, Inches(0.35),
                 font_size=Pt(20), color=C_WHITE)
        y += Inches(0.4)
    add_text(slide, title, left+pad, y, width-2*pad, Inches(0.35),
             font_size=Pt(13), bold=True, color=title_color)
    y += Inches(0.38)
    add_text(slide, body, left+pad, y, width-2*pad, height - (y - top) - pad,
             font_size=Pt(11), color=body_color, word_wrap=True)


def divider(slide, top, color=C_PURPLE):
    add_rect(slide, Inches(0.5), top, W - Inches(1), Pt(1),
             fill_color=color, line_color=color)


def bullet_item(slide, text, left, top, width, checked=True, check_color=C_GREEN):
    mark = "✓" if checked else "✗"
    col  = check_color if checked else C_RED
    add_text(slide, mark, left, top, Inches(0.3), Inches(0.3),
             font_size=Pt(13), bold=True, color=col)
    add_text(slide, text, left + Inches(0.35), top, width - Inches(0.35), Inches(0.3),
             font_size=Pt(12), color=C_MUTED)


def status_row(slide, label, left, top, width):
    # green dot
    dot = add_rect(slide, left, top + Inches(0.06), Inches(0.14), Inches(0.14),
                   fill_color=C_GREEN, line_color=C_GREEN)
    add_text(slide, label, left + Inches(0.22), top, width - Inches(0.22), Inches(0.28),
             font_size=Pt(12), color=C_WHITE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — HERO
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)

add_chip(s, "🎓  Startup 101  ·  MVP v3  ·  April 2026",
         h_center(Inches(4)), Inches(1.1))

add_text(s, "CampusConnect", h_center(Inches(10)), Inches(1.7),
         Inches(10), Inches(1.6), font_size=Pt(72), bold=True,
         color=C_VIOLET, align=PP_ALIGN.CENTER)

add_text(s,
         "A peer-to-peer anonymous digital marketplace &\nAI-powered Lost & Found — built for university campuses.",
         h_center(Inches(9)), Inches(3.4), Inches(9), Inches(0.9),
         font_size=Pt(18), color=C_MUTED, align=PP_ALIGN.CENTER)

tags = ["🛒 Marketplace", "🔍 AI Lost & Found", "🛡️ Role-Based", "💳 Stripe", "👤 Anonymous"]
tag_w = Inches(1.9)
total_tag_w = tag_w * len(tags) + Inches(0.2) * (len(tags)-1)
tx = h_center(total_tag_w)
for t in tags:
    r = add_rect(s, tx, Inches(4.6), tag_w, Inches(0.45),
                 fill_color=RGBColor(0x20,0x20,0x38),
                 line_color=RGBColor(0x44,0x44,0x66), line_width=Pt(1))
    tf = r.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = t
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    tx += tag_w + Inches(0.2)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x1A, 0x0A, 0x0A))

add_chip(s, "🔥  The Problem", h_center(Inches(2.5)), Inches(0.3))
add_text(s, "Campus trading is broken", h_center(Inches(10)), Inches(0.8),
         Inches(10), Inches(0.8), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

pains = [
    ("📱", "WhatsApp Chaos",        "Lost & sale messages buried within hours. No structure, no history."),
    ("🎭", "No Privacy",            "Real identities exposed in every transaction — safety & trust concern."),
    ("💸", "No Payment Safety",     "Cash or UPI only — no verification, high scam risk for students."),
    ("🔎", "Lost Objects Never Recovered", "Manual notice boards are slow, unreliable, and not searchable."),
]
py = Inches(1.8)
for icon, title, body in pains:
    add_text(s, icon,  Inches(0.5), py, Inches(0.5), Inches(0.4), font_size=Pt(20))
    add_text(s, title, Inches(1.1), py, Inches(5.0), Inches(0.3), font_size=Pt(13), bold=True, color=C_WHITE)
    add_text(s, body,  Inches(1.1), py+Inches(0.32), Inches(5.4), Inches(0.35), font_size=Pt(11), color=C_MUTED)
    py += Inches(0.85)

# quote box right side
qx, qy, qw, qh = Inches(7.0), Inches(1.8), Inches(5.8), Inches(1.4)
qr = add_rect(s, qx, qy, qw, qh, fill_color=RGBColor(0x20,0x10,0x10),
              line_color=C_PURPLE, line_width=Pt(2))
tf = qr.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = '"When someone\'s hair is on fire,\nthey\'ll take a brick if it helps."\n— Paul Graham, YC'
run.font.size  = Pt(13)
run.font.color.rgb = C_MUTED
run.font.italic = True

scenarios = [
    "😰  'I lost my wallet before exams and can't reach anyone'",
    "📦  'I'm leaving campus Friday — need to sell everything NOW'",
]
sy = Inches(3.5)
for sc in scenarios:
    add_text(s, sc, Inches(7.0), sy, Inches(5.8), Inches(0.4), font_size=Pt(12), color=RGBColor(0xF8,0x71,0x71))
    sy += Inches(0.5)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — SOLUTION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x03, 0x1A, 0x1A))

add_chip(s, "💡  Our Solution", h_center(Inches(2.5)), Inches(0.3))
add_text(s, "One platform. Everything campus.", h_center(Inches(11)), Inches(0.85),
         Inches(11), Inches(0.75), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

cols = [
    ("🛒", "Anonymous Marketplace",
     "Buy & sell items with unique anonymous handles.\nNeither buyer nor seller knows each other's real identity."),
    ("🔍", "AI-Powered Lost & Found",
     "TF-IDF text + OpenCV image matching auto-connects lost reports with found items in real-time."),
    ("🛡️", "Admin Dashboard",
     "Real-time analytics, revenue tracking, dispute resolution — all isolated from marketplace activity."),
]
cw = Inches(3.8)
cx = Inches(0.65)
for icon, title, body in cols:
    card(s, cx, Inches(2.0), cw, Inches(3.8), title, body, icon=icon,
         border=C_PURPLE)
    cx += cw + Inches(0.35)

qr = add_rect(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.85),
              fill_color=RGBColor(0x10,0x20,0x28),
              line_color=C_PURPLE, line_width=Pt(1.5))
tf = qr.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '"Unlike Rapido & Zomato — CampusConnect is two-way anonymous. The platform itself is the trusted third party."'
run.font.size = Pt(12); run.font.color.rgb = C_MUTED; run.font.italic = True

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — TECH STACK
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x0A, 0x0A, 0x1A))

add_chip(s, "🏗️  Architecture", h_center(Inches(2.5)), Inches(0.3))
add_text(s, "Full-Stack Tech Stack", h_center(Inches(9)), Inches(0.85),
         Inches(9), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

# Left: architecture diagram
arch = [
    ("⚛️  Frontend — React.js + Vite + Tailwind CSS", C_INDIGO),
    ("         ↕  Axios REST API", C_MUTED),
    ("🐍  Backend — Flask + JWT Auth + Gunicorn", C_PURPLE),
    ("         ↕  ↕", C_MUTED),
    ("🍃  MongoDB Atlas", C_GREEN),
    ("💳  Stripe API", C_BLUE),
]
ay = Inches(1.8)
for line, col in arch:
    add_text(s, line, Inches(0.5), ay, Inches(5.5), Inches(0.4),
             font_size=Pt(13), color=col, bold=(col != C_MUTED))
    ay += Inches(0.55)

# Right: table
rows = [
    ("Frontend",  "React, Vite, Tailwind CSS"),
    ("Backend",   "Flask, JWT, Bcrypt"),
    ("Database",  "MongoDB Atlas"),
    ("Auth",      "JWT + Google OAuth 2.0"),
    ("Payments",  "Stripe Checkout API"),
    ("AI Engine", "TF-IDF + OpenCV"),
    ("Deploy",    "Gunicorn + Render"),
]
ty, col1_x, col2_x = Inches(1.9), Inches(7.0), Inches(9.6)
hdr = add_rect(s, col1_x, ty, Inches(5.8), Inches(0.4),
               fill_color=RGBColor(0x2A,0x1A,0x55))
add_text(s, "Layer",      col1_x+Inches(0.1), ty+Inches(0.05), Inches(2.4), Inches(0.3), font_size=Pt(12), bold=True, color=C_VIOLET)
add_text(s, "Technology", col2_x+Inches(0.1), ty+Inches(0.05), Inches(3.2), Inches(0.3), font_size=Pt(12), bold=True, color=C_VIOLET)
ty += Inches(0.42)
for layer, tech in rows:
    br = add_rect(s, col1_x, ty, Inches(5.8), Inches(0.38),
                  fill_color=RGBColor(0x18,0x18,0x2E),
                  line_color=RGBColor(0x30,0x30,0x50), line_width=Pt(0.5))
    add_text(s, layer, col1_x+Inches(0.1), ty+Inches(0.04), Inches(2.4), Inches(0.3), font_size=Pt(11), color=C_MUTED)
    add_text(s, tech,  col2_x+Inches(0.1), ty+Inches(0.04), Inches(3.2), Inches(0.3), font_size=Pt(11), bold=True, color=C_WHITE)
    ty += Inches(0.38)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — USER ROLES
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x0A, 0x1A, 0x0A))

add_chip(s, "👥  User Roles", h_center(Inches(2.5)), Inches(0.3))
add_text(s, "3 Distinct Roles", h_center(Inches(8)), Inches(0.85),
         Inches(8), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

roles = [
    ("🎓", "Student", C_INDIGO, [
        (True,  "Login via Email or Google"),
        (True,  "List & Buy marketplace items"),
        (True,  "Report Lost & Claim Found"),
        (True,  "Anonymous handle (User#ABC12)"),
        (True,  "Secure Stripe checkout"),
    ]),
    ("🏫", "Alumni", C_GREEN, [
        (True,  "Same access as Student"),
        (True,  "Separate role badge"),
        (True,  "Access via email/Google"),
        (True,  "Anonymous identity enforced"),
        (True,  "All marketplace features"),
    ]),
    ("⚙️", "Admin", C_RED, [
        (True,  "Full dashboard access"),
        (True,  "Sees real names for disputes"),
        (True,  "Revenue & analytics view"),
        (False, "Cannot buy/sell/claim"),
        (False, "Platform neutrality enforced"),
    ]),
]
rw = Inches(3.8)
rx = Inches(0.65)
for icon, title, border, items in roles:
    # card bg
    add_rect(s, rx, Inches(1.9), rw, Inches(4.6),
             fill_color=RGBColor(0x1A,0x1A,0x2E),
             line_color=border, line_width=Pt(1.5))
    add_text(s, icon,  rx+Inches(0.2), Inches(2.05), Inches(0.5), Inches(0.4), font_size=Pt(24))
    add_text(s, title, rx+Inches(0.75), Inches(2.05), rw-Inches(0.9), Inches(0.4),
             font_size=Pt(16), bold=True, color=C_WHITE)
    iy = Inches(2.65)
    for checked, txt in items:
        bullet_item(s, txt, rx+Inches(0.2), iy, rw-Inches(0.35), checked)
        iy += Inches(0.45)
    rx += rw + Inches(0.35)

# bottom quote
qr = add_rect(s, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.55),
              fill_color=RGBColor(0x10,0x18,0x10), line_color=C_PURPLE, line_width=Pt(1))
tf = qr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '"Admin isolation enforced on both frontend AND every backend route — double security layer."'
run.font.size = Pt(11); run.font.color.rgb = C_MUTED; run.font.italic = True

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — MARKETPLACE FLOW
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x1A, 0x0F, 0x00))

add_chip(s, "🛒  Marketplace", h_center(Inches(2.5)), Inches(0.3))
add_text(s, "How the Marketplace Works", h_center(Inches(10)), Inches(0.85),
         Inches(10), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

seller_steps = [
    "1.  Enter item name, category, condition & price",
    "2.  Upload image (stored securely)",
    "3.  Item listed with anonymous handle",
    "4.  Seller can edit / delete anytime",
]
buyer_steps  = [
    "1.  Browse, search, filter by price & category",
    "2.  Click 'Buy Now' → see fee breakdown",
    "3.  Redirect to Stripe Checkout",
    "4.  Item marked 'Sold' → Admin updated",
]

add_text(s, "SELLER FLOW", Inches(0.5), Inches(1.8), Inches(5), Inches(0.3),
         font_size=Pt(11), bold=True, color=C_ORANGE)
sy = Inches(2.2)
for step in seller_steps:
    add_text(s, step, Inches(0.5), sy, Inches(5.5), Inches(0.38), font_size=Pt(12), color=C_MUTED)
    sy += Inches(0.48)

add_text(s, "BUYER FLOW", Inches(6.2), Inches(1.8), Inches(5), Inches(0.3),
         font_size=Pt(11), bold=True, color=C_BLUE)
by = Inches(2.2)
for step in buyer_steps:
    add_text(s, step, Inches(6.2), by, Inches(5.5), Inches(0.38), font_size=Pt(12), color=C_MUTED)
    by += Inches(0.48)

# Benefits card
benef = add_rect(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8),
                 fill_color=RGBColor(0x18,0x18,0x2E),
                 line_color=RGBColor(0x44,0x44,0x66), line_width=Pt(1))
bx = Inches(0.9)
for icon, title, desc in [("💼", "Seller", "Gets full listed price"),
                           ("🛒", "Buyer",  "Gets verified item + safe tx"),
                           ("🏛️", "Platform","Earns % fee per sale")]:
    add_text(s, icon,  bx, Inches(5.0), Inches(0.5), Inches(0.5), font_size=Pt(24))
    add_text(s, title, bx+Inches(0.55), Inches(5.0), Inches(3.5), Inches(0.3),
             font_size=Pt(13), bold=True, color=C_WHITE)
    add_text(s, desc,  bx+Inches(0.55), Inches(5.4), Inches(3.5), Inches(0.3),
             font_size=Pt(11), color=C_MUTED)
    bx += Inches(4.2)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — BUSINESS MODEL
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x0F, 0x0A, 0x1A))

add_chip(s, "💰  Business Model", h_center(Inches(2.8)), Inches(0.3))
add_text(s, "Tiered Fee Structure", h_center(Inches(9)), Inches(0.85),
         Inches(9), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

fees = [
    ("📚  Books",        "3%",   C_GREEN,  0.25),
    ("💻  Electronics",  "12%",  C_RED,    1.00),
    ("₹0 – ₹2,000",     "5%",   C_BLUE,   0.42),
    ("₹2,001 – ₹5,000", "8%",   C_VIOLET, 0.67),
    ("₹5,001+",         "10%",  C_ORANGE, 0.83),
]
fy, bar_x, bar_w = Inches(2.0), Inches(3.6), Inches(6.0)
add_text(s, "CATEGORY-BASED (Priority)", Inches(0.5), fy-Inches(0.3), Inches(8), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_MUTED)
for i, (label, pct, col, ratio) in enumerate(fees):
    if i == 2:
        fy += Inches(0.3)
        add_text(s, "PRICE-BASED (Other Categories)", Inches(0.5), fy-Inches(0.05), Inches(8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=C_MUTED)
        fy += Inches(0.3)
    add_text(s, label, Inches(0.5), fy, Inches(2.8), Inches(0.3), font_size=Pt(12), color=C_MUTED)
    bg_bar = add_rect(s, bar_x, fy+Inches(0.04), bar_w, Inches(0.2),
                      fill_color=RGBColor(0x25,0x25,0x40))
    if ratio > 0:
        add_rect(s, bar_x, fy+Inches(0.04), Inches(bar_w.inches * ratio), Inches(0.2),
                 fill_color=col)
    add_text(s, pct, bar_x+bar_w+Inches(0.1), fy, Inches(0.7), Inches(0.3),
             font_size=Pt(12), bold=True, color=col)
    fy += Inches(0.55)

# Right: revenue example
re_x = Inches(9.9)
add_rect(s, re_x, Inches(2.0), Inches(3.2), Inches(3.8),
         fill_color=RGBColor(0x18,0x18,0x2E), line_color=C_PURPLE, line_width=Pt(1))
add_text(s, "REVENUE EXAMPLE", re_x+Inches(0.15), Inches(2.15), Inches(3.0), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_VIOLET)
examples = [
    ("Book ₹300",      "+₹9",    C_GREEN),
    ("Laptop ₹25,000", "+₹3,000",C_RED),
    ("Cycle ₹5,500",   "+₹550",  C_ORANGE),
    ("Table ₹1,800",   "+₹90",   C_BLUE),
]
ey = Inches(2.6)
for item, rev, col in examples:
    add_text(s, item, re_x+Inches(0.15), ey, Inches(1.9), Inches(0.3), font_size=Pt(11), color=C_MUTED)
    add_text(s, rev,  re_x+Inches(2.2), ey, Inches(0.9), Inches(0.3), font_size=Pt(12), bold=True, color=col)
    ey += Inches(0.5)

# L&F free badge
add_rect(s, re_x, Inches(6.0), Inches(3.2), Inches(0.9),
         fill_color=RGBColor(0x0A,0x25,0x15), line_color=C_GREEN, line_width=Pt(1.5))
add_text(s, "🔍  Lost & Found: FREE (0%)", re_x+Inches(0.1), Inches(6.1), Inches(3.0), Inches(0.35),
         font_size=Pt(12), bold=True, color=C_GREEN)
add_text(s, "Builds community trust first", re_x+Inches(0.1), Inches(6.45), Inches(3.0), Inches(0.3),
         font_size=Pt(10), color=C_MUTED)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — AI LOST & FOUND
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x1A, 0x0A, 0x0F))

add_chip(s, "🤖  AI Lost & Found", h_center(Inches(3)), Inches(0.3))
add_text(s, "Smart Matching Engine", h_center(Inches(9)), Inches(0.85),
         Inches(9), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

ai_cards = [
    ("📝", "TF-IDF Text Similarity",  "Compares descriptions word-by-word using cosine similarity."),
    ("🖼️", "OpenCV Image Hashing",    "Perceptual hash detects visually similar items, even edited ones."),
    ("🎯", "Category & Name Bonus",   "+15% for same category; +20% if item names are substrings."),
]
cy = Inches(2.0)
for icon, title, body in ai_cards:
    card(s, Inches(0.4), cy, Inches(5.2), Inches(1.35), title, body, icon=icon, border=C_PINK)
    cy += Inches(1.5)

# Formula box
fb = add_rect(s, Inches(6.2), Inches(2.0), Inches(6.8), Inches(3.2),
              fill_color=RGBColor(0x1A,0x08,0x12), line_color=C_PINK, line_width=Pt(1.5))
add_text(s, "COMBINED SCORE FORMULA", Inches(6.4), Inches(2.15), Inches(6.4), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_PINK)
formula = (
    "With Images:\n"
    "  Score = Text×0.40 + Image×0.35 + Bonus×0.25\n\n"
    "Without Images:\n"
    "  Score = Text×0.70 + Bonus×0.30\n\n"
    "MATCH if Score > 0.25"
)
add_text(s, formula, Inches(6.4), Inches(2.55), Inches(6.4), Inches(2.5),
         font_size=Pt(12), color=C_VIOLET)

# Claim steps
add_text(s, "CLAIM PROCESS", Inches(0.4), Inches(5.65), Inches(5.5), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_PINK)
claim_steps = [
    "1️⃣  Report lost → AI scans all 'Found' reports",
    "2️⃣  Matched items shown ranked by score",
    "3️⃣  Claimer submits proof of ownership",
    "4️⃣  Admin verifies & approves claim",
    "5️⃣  Item marked Resolved ✓",
]
csy = Inches(6.05)
for step in claim_steps:
    add_text(s, step, Inches(0.4), csy, Inches(12.5), Inches(0.3), font_size=Pt(11), color=C_MUTED)
    csy += Inches(0.4)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — ADMIN DASHBOARD
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x00, 0x1A, 0x2A))

add_chip(s, "📊  Admin Dashboard", h_center(Inches(3)), Inches(0.3))
add_text(s, "Platform Control Room", h_center(Inches(9)), Inches(0.85),
         Inches(9), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

admin_items = [
    ("👥", "Total Users",       "All registered Students & Alumni",    C_YELLOW),
    ("💰", "Platform Revenue",  "Net earnings from all marketplace fees",C_GREEN),
    ("📦", "Items Sold",        "Complete transaction history log",     C_BLUE),
    ("🔎", "Lost Reports",      "Open & Resolved lost report counts",   C_PINK),
    ("✅", "Found Reports",     "All submitted found item reports",     C_ORANGE),
    ("🤝", "Claims",            "Pending / Approved / Rejected",        C_VIOLET),
    ("📋", "Transactions",      "Full buyer-seller payment ledger",     RGBColor(0x34,0xD3,0x99)),
    ("🛡️", "Role Isolated",    "Cannot buy, sell, or claim items",     C_RED),
]
cols_n = 4
cw, ch = Inches(3.0), Inches(1.5)
gap = Inches(0.2)
total_w = cols_n * cw + (cols_n-1) * gap
sx = h_center(total_w)
sy = Inches(1.9)
for i, (icon, title, body, col) in enumerate(admin_items):
    cx = sx + (i % cols_n) * (cw + gap)
    cy = sy + (i // cols_n) * (ch + gap)
    card(s, cx, cy, cw, ch, title, body, icon=icon, border=col)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — MVP ITERATIONS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x00, 0x1A, 0x0A))

add_chip(s, "🔄  Build → Measure → Learn", h_center(Inches(4)), Inches(0.3))
add_text(s, "MVP Iterations & Pivots", h_center(Inches(9)), Inches(0.85),
         Inches(9), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

# Table
iter_rows = [
    ("v1", C_BLUE,   "Flat 5% fee on all items",      "Users didn't differentiate",       "Kept as baseline"),
    ("v2", C_ORANGE, "L&F claim fee (variable)",       "Users chose free WhatsApp instead","Fee was blocking adoption"),
    ("v3", C_GREEN,  "L&F = FREE (0%)",                "Community trust metric",           "→ Removed fee entirely ✓"),
    ("v3", C_GREEN,  "Category + Tiered pricing",      "Electronics = higher margin",      "→ Books 3%, Elec 12% ✓"),
]
headers = ["Ver", "Hypothesis", "Insight", "Pivot"]
hh = Inches(0.4)
hw = [Inches(0.6), Inches(3.2), Inches(3.2), Inches(3.2)]
hx = Inches(0.5)
ty = Inches(1.9)

# Header row
hbg = add_rect(s, Inches(0.5), ty, sum(hw)+Inches(0.15), hh,
               fill_color=RGBColor(0x1A,0x40,0x28))
cx = Inches(0.6)
for h_label, hw_ in zip(headers, hw):
    add_text(s, h_label, cx, ty+Inches(0.05), hw_, Inches(0.3),
             font_size=Pt(11), bold=True, color=C_VIOLET)
    cx += hw_
ty += hh

for ver, col, hyp, insight, pivot in iter_rows:
    rbg = add_rect(s, Inches(0.5), ty, sum(hw)+Inches(0.15), hh,
                   fill_color=RGBColor(0x14,0x24,0x18),
                   line_color=RGBColor(0x20,0x40,0x28), line_width=Pt(0.5))
    cx = Inches(0.6)
    add_text(s, ver,     cx, ty+Inches(0.05), hw[0], Inches(0.3), font_size=Pt(11), bold=True, color=col)
    cx += hw[0]
    add_text(s, hyp,     cx, ty+Inches(0.05), hw[1], Inches(0.3), font_size=Pt(11), bold=True, color=C_WHITE)
    cx += hw[1]
    add_text(s, insight, cx, ty+Inches(0.05), hw[2], Inches(0.3), font_size=Pt(11), color=C_MUTED)
    cx += hw[2]
    add_text(s, pivot,   cx, ty+Inches(0.05), hw[3], Inches(0.3), font_size=Pt(11), color=C_GREEN)
    ty += hh

# Right column: quotes + FR8 insight
qx2 = Inches(10.8)
add_rect(s, qx2, Inches(1.9), Inches(2.3), Inches(1.5),
         fill_color=RGBColor(0x10,0x20,0x10), line_color=C_PURPLE, line_width=Pt(1))
add_text(s,
         '"If you\'re not embarrassed by the first version, you\'ve launched too late."\n— Reid Hoffman',
         qx2+Inches(0.1), Inches(2.0), Inches(2.1), Inches(1.3),
         font_size=Pt(10), color=C_MUTED, italic=True)

add_rect(s, qx2, Inches(3.55), Inches(2.3), Inches(2.0),
         fill_color=RGBColor(0x08,0x20,0x10), line_color=C_GREEN, line_width=Pt(1))
add_text(s, "YC PRINCIPLE APPLIED ✓",
         qx2+Inches(0.1), Inches(3.65), Inches(2.1), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_GREEN)
add_text(s,
         "Like FR8 — we identified information asymmetry. WhatsApp = 30%+ lost deals. "
         "CampusConnect is the broker-free campus orchestration layer.",
         qx2+Inches(0.1), Inches(4.0), Inches(2.1), Inches(1.3),
         font_size=Pt(10), color=C_MUTED, word_wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — STATUS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0x0F, 0x0A, 0x1A))

add_chip(s, "✅  Platform Status", h_center(Inches(3)), Inches(0.3))
add_text(s, "MVP v3 is Production Ready", h_center(Inches(10)), Inches(0.85),
         Inches(10), Inches(0.7), font_size=Pt(40), bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

live = [
    "Email + Google OAuth Login",
    "Anonymous Handle Assignment",
    "Marketplace (List / Browse / Buy)",
    "Stripe Payment Integration",
    "Tiered + Category Fee Engine",
    "Lost & Found Reporting (FREE)",
]
also = [
    "AI Smart Matching (Text + Image)",
    "Verified Claim Workflow",
    "Admin Dashboard (Live Analytics)",
    "Admin Role Isolation (Backend+UI)",
    "Password Reset via Email",
    "Gunicorn Production Deploy",
]
sy2 = Inches(2.0)
for item in live:
    status_row(s, item, Inches(0.6), sy2, Inches(6.0))
    sy2 += Inches(0.55)

sy2 = Inches(2.0)
for item in also:
    status_row(s, item, Inches(7.0), sy2, Inches(6.0))
    sy2 += Inches(0.55)

# Roadmap tags
add_text(s, "UPCOMING ROADMAP", Inches(0.6), Inches(6.3), Inches(6), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_MUTED)
roadmap = ["🔮 In-App Chat", "⭐ Ratings", "📱 Mobile App", "🔔 Push Notifs", "🏫 Partnerships"]
rx2 = Inches(0.6)
for item in roadmap:
    tw = Inches(len(item) * 0.09 + 0.4)
    rr = add_rect(s, rx2, Inches(6.65), tw, Inches(0.42),
                  fill_color=RGBColor(0x1E,0x1E,0x35),
                  line_color=RGBColor(0x44,0x44,0x66), line_width=Pt(1))
    tf = rr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = item
    run.font.size = Pt(11); run.font.color.rgb = C_MUTED; run.font.bold = True
    rx2 += tw + Inches(0.15)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CLOSING
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)

add_chip(s, "🎓  Thank You", h_center(Inches(2.5)), Inches(0.8))
add_text(s, "CampusConnect", h_center(Inches(10)), Inches(1.4),
         Inches(10), Inches(1.3), font_size=Pt(72), bold=True,
         color=C_VIOLET, align=PP_ALIGN.CENTER)
add_text(s, "Built with the Y Combinator philosophy:",
         h_center(Inches(9)), Inches(2.9), Inches(9), Inches(0.4),
         font_size=Pt(16), color=C_MUTED, align=PP_ALIGN.CENTER)
add_text(s, "Launch fast. Iterate constantly. Listen to users.",
         h_center(Inches(9)), Inches(3.3), Inches(9), Inches(0.5),
         font_size=Pt(20), bold=True, color=C_VIOLET, align=PP_ALIGN.CENTER)

icons_end = [("🛒", "Peer-to-Peer\nMarketplace"),
             ("🤖", "AI Lost &\nFound"),
             ("💰", "Smart Tiered\nBusiness Model")]
ix = Inches(2.0)
for icon, label in icons_end:
    add_text(s, icon, ix, Inches(4.1), Inches(3.0), Inches(0.6), font_size=Pt(36), align=PP_ALIGN.CENTER)
    add_text(s, label, ix, Inches(4.8), Inches(3.0), Inches(0.7), font_size=Pt(13), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    ix += Inches(3.4)

footer = add_rect(s, h_center(Inches(7)), Inches(5.9), Inches(7), Inches(0.6),
                  fill_color=RGBColor(0x18,0x10,0x2E),
                  line_color=C_PURPLE, line_width=Pt(1))
tf = footer.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Built at IIIT Sricity  ·  MVP v3  ·  Startup 101  ·  April 2026"
run.font.size = Pt(12); run.font.color.rgb = C_MUTED

# ─── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\amitk\OneDrive\Desktop\FullStack Projects\campus_connecting\campus-connecting\CampusConnect_Presentation.pptx"
prs.save(out)
print(f"[DONE] Saved -> {out}")
